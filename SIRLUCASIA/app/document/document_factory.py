import json
from docx import Document
from reportlab.pdfgen import canvas
from openpyxl import Workbook

try:
    from pptx import Presentation
except ImportError as e:
    raise RuntimeError("python-pptx no está instalado en este entorno") from e


class DocumentFactory:

    # =====================================
    # TXT
    # =====================================
    @staticmethod
    def create_txt(path, content):
        with open(path, "w", encoding="utf-8") as file:
            file.write(content)

    # =====================================
    # Markdown
    # =====================================
    @staticmethod
    def create_md(path, content):
        with open(path, "w", encoding="utf-8") as file:
            file.write(content)

    # =====================================
    # JSON
    # =====================================
    @staticmethod
    def create_json(path, content):
        try:
            data = json.loads(content)
            with open(path, "w", encoding="utf-8") as file:
                json.dump(data, file, indent=4, ensure_ascii=False)
        except Exception:
            with open(path, "w", encoding="utf-8") as file:
                file.write(content)

    # =====================================
    # PDF
    # =====================================
    @staticmethod
    def create_pdf(path, content):
        pdf = canvas.Canvas(path)
        y = 800
        for line in content.split("\n"):
            pdf.drawString(50, y, line)
            y -= 20
        pdf.save()

    # =====================================
    # Word
    # =====================================
    @staticmethod
    def create_docx(path, content):
        document = Document()
        document.add_heading("Documento generado por SIRLUCAS AI", level=1)
        for line in content.split("\n"):
            document.add_paragraph(line)
        document.save(path)

    # =====================================
    # Excel
    # =====================================
    @staticmethod
    def create_xlsx(path, content):
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Datos"
        for row, line in enumerate(content.split("\n"), start=1):
            sheet.cell(row=row, column=1).value = line
        workbook.save(path)

    # =====================================
    # PowerPoint
    # =====================================
    @staticmethod
    def create_pptx(path, content):
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # título + contenido
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "Documento generado por SIRLUCAS AI"
        body.text = content
        prs.save(path)

    # =====================================
    # Dispatcher
    # =====================================
    @staticmethod
    def create(extension, path, content):
        creators = {
            ".txt": DocumentFactory.create_txt,
            ".md": DocumentFactory.create_md,
            ".json": DocumentFactory.create_json,
            ".pdf": DocumentFactory.create_pdf,
            ".docx": DocumentFactory.create_docx,
            ".xlsx": DocumentFactory.create_xlsx,
            ".pptx": DocumentFactory.create_pptx
        }

        creator = creators.get(extension)
        if creator is None:
            raise ValueError(f"Formato no soportado: {extension}")

        creator(path, content)
